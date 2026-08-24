"""Generate a colorful, modern PowerPoint presentation documenting the automation architecture.

Usage:
    python scripts/generate_pptx.py --palette blue_teal --output docs/automation_presentation_blue_teal.pptx
    python scripts/generate_pptx.py --palette purple_orange --output docs/automation_presentation_purple_orange.pptx
    python scripts/generate_pptx.py --palette both

This script creates PPTX files but does NOT commit them to the repository. Run locally and then decide whether to commit the binary.
"""
import argparse
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
import os


def set_background_color(slide, rgb):
    # Add a full-size rectangle and fill with the color to emulate background change
    left = top = Inches(0)
    width = Inches(10)
    height = Inches(7.5)
    shape = slide.shapes.add_shape(1, left, top, width, height)  # MSO_SHAPE.RECTANGLE = 1
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(*rgb)
    shape.line.fill.background()  # remove border
    return shape


def add_title_slide(prs, title_text, subtitle_text, palette):
    slide_layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    set_background_color(slide, palette['bg'])

    # Title
    left = Inches(0.8)
    top = Inches(1.0)
    width = Inches(8.5)
    height = Inches(1.5)
    title_box = slide.shapes.add_textbox(left, top, width, height)
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(*palette['title'])

    # Subtitle
    left = Inches(0.8)
    top = Inches(2.4)
    width = Inches(8.5)
    height = Inches(1.0)
    sub_box = slide.shapes.add_textbox(left, top, width, height)
    st = sub_box.text_frame
    p2 = st.paragraphs[0]
    p2.text = subtitle_text
    p2.font.size = Pt(14)
    p2.font.color.rgb = RGBColor(*palette['subtitle'])


def add_bulleted_slide(prs, title, bullets, palette):
    slide_layout = prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    # light background
    set_background_color(slide, palette['light_bg'])

    # Title
    title_shape = slide.shapes.title
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.size = Pt(28)
    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(*palette['title'])

    # Body
    body = slide.shapes.placeholders[1].text_frame
    body.clear()
    for i, b in enumerate(bullets):
        if i == 0:
            p = body.paragraphs[0]
            p.text = b
        else:
            p = body.add_paragraph()
            p.text = b
        p.level = 0
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(*palette['text'])


def build_presentation_for_palette(output_path, palette_name):
    palettes = {
        'blue_teal': {
            'bg': (4, 77, 111),         # deep blue
            'title': (255, 255, 255),   # white
            'subtitle': (200, 230, 240),
            'light_bg': (230, 245, 250),
            'text': (20, 40, 60)
        },
        'purple_orange': {
            'bg': (80, 24, 120),        # deep purple
            'title': (255, 255, 255),
            'subtitle': (255, 230, 200),
            'light_bg': (250, 240, 235),
            'text': (50, 30, 30)
        }
    }

    if palette_name not in palettes:
        raise ValueError('Unknown palette: ' + palette_name)

    palette = palettes[palette_name]
    prs = Presentation()

    # Title
    add_title_slide(prs, 'Python Selenium Automation Boilerplate', 'Architecture, files and how-to guide', palette)

    # Objectives
    add_bulleted_slide(prs, 'Objectives', [
        'Reusable baseline for web UI automation',
        'Page Object Model (POM)',
        'pytest-based tests and reporting',
        'Utilities & assertions'
    ], palette)

    # Repo structure
    add_bulleted_slide(prs, 'Repository Structure', [
        'pages/: Page objects (BasePage, HomePage, LoginPage)',
        'tests/: pytest tests and fixtures',
        'utils/: driver, config, logger, screenshots',
        'asserts/: assertion helpers'
    ], palette)

    # POM concept
    add_bulleted_slide(prs, 'Page Object Model (POM)', [
        'One class per page with locators & actions',
        'Keeps tests readable and robust to UI changes',
        'Example: HomePage.open(), LoginPage.login()'
    ], palette)

    # Tests design
    add_bulleted_slide(prs, 'Tests & Fixtures', [
        'Use pytest fixtures for driver lifecycle',
        'Group tests with markers (smoke, regression)',
        'Keep tests independent for parallel runs'
    ], palette)

    # Utilities & assertions
    add_bulleted_slide(prs, 'Utilities & Assertions', [
        'utils/: driver factory, config reader, screenshots',
        'asserts/: centralized assertion helpers with screenshots',
        'reporting: pytest-html and artifacts in CI'
    ], palette)

    # Config & Running
    add_bulleted_slide(prs, 'Configuration & Running', [
        'config.ini for base_url, browser, timeouts',
        'BROWSER & HEADLESS env vars override config',
        'pip install -r requirements.txt',
        'pytest -v --html=reports/python_html_report.html'
    ], palette)

    # CI notes
    add_bulleted_slide(prs, 'CI Notes (Jenkins / GitHub Actions)', [
        'Checkout repo, setup Python, pip install -r requirements.txt',
        'Run pytest and archive reports/screenshots',
        'Use secrets store for credentials (Jenkins/GitHub Secrets)'
    ], palette)

    # Best practices
    add_bulleted_slide(prs, 'Best Practices', [
        'Use explicit waits; avoid sleeps',
        'Centralize assertions for consistent messages & screenshots',
        'Parallelize tests with pytest-xdist; keep tests isolated'
    ], palette)

    # Next steps
    add_bulleted_slide(prs, 'Next steps', [
        'Add more page objects & integration tests',
        'Add GitHub Actions CI matrix',
        'Consider Dockerized test agents or cloud grids (BrowserStack)'
    ], palette)

    # Save
    os.makedirs(os.path.dirname(output_path), exist_ok=True)
    prs.save(output_path)
    print(f'Presentation written: {output_path}')


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument('--palette', default='both', choices=['blue_teal', 'purple_orange', 'both'])
    parser.add_argument('--output', default='docs/automation_presentation.pptx')
    args = parser.parse_args()

    if args.palette == 'both':
        out1 = os.path.splitext(args.output)[0] + '_blue_teal.pptx'
        out2 = os.path.splitext(args.output)[0] + '_purple_orange.pptx'
        build_presentation_for_palette(out1, 'blue_teal')
        build_presentation_for_palette(out2, 'purple_orange')
    else:
        build_presentation_for_palette(args.output, args.palette)


if __name__ == '__main__':
    main()
